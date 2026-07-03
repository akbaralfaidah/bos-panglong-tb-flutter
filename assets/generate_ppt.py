"""
Generate Bos Depot Apps Portfolio PPT
Clean White + Navy Blue Theme | Poppins Font
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

# ─── CONSTANTS ───
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY       = RGBColor(0x0A, 0x26, 0x47)
NAVY_LIGHT = RGBColor(0x14, 0x42, 0x72)
GOLD       = RGBColor(0xF3, 0x9C, 0x12)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT   = RGBColor(0xF4, 0xF7, 0xF6)
TEXT_DARK   = RGBColor(0x1E, 0x29, 0x3B)
TEXT_GREY   = RGBColor(0x64, 0x74, 0x8B)
GREEN      = RGBColor(0x19, 0x87, 0x54)
RED        = RGBColor(0xDC, 0x35, 0x45)
TEAL       = RGBColor(0x00, 0x79, 0x6B)
TEAL_BG    = RGBColor(0xE0, 0xF2, 0xF1)
BLUE_BG    = RGBColor(0xE3, 0xF2, 0xFD)
AMBER_BG   = RGBColor(0xFF, 0xF8, 0xE1)
GREEN_BG   = RGBColor(0xE8, 0xF5, 0xE9)
PINK_BG    = RGBColor(0xFC, 0xE4, 0xEC)
INDIGO_BG  = RGBColor(0xE8, 0xEA, 0xF6)
SURFACE    = RGBColor(0xEB, 0xF4, 0xFF)

FONT = "Poppins"
IMG_DIR = os.path.join(os.path.dirname(__file__), "images")
OUT_FILE = os.path.join(os.path.dirname(__file__), "Bos_Depot_Apps_Portfolio.pptx")


# ─── HELPERS ───
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def set_slide_gradient(slide, c1, c2):
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = c1
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = c2
    fill.gradient_stops[1].position = 1.0

def add_text_box(slide, left, top, width, height, text, font_size=18, color=TEXT_DARK,
                  bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_multiline_box(slide, left, top, width, height, lines, default_size=14,
                       default_color=TEXT_DARK, alignment=PP_ALIGN.LEFT):
    """lines = list of dicts: {text, size, color, bold, spacing_after}"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line.get("text", "")
        p.font.size = Pt(line.get("size", default_size))
        p.font.color.rgb = line.get("color", default_color)
        p.font.bold = line.get("bold", False)
        p.font.name = FONT
        p.alignment = alignment
        p.space_after = Pt(line.get("spacing_after", 4))
        p.space_before = Pt(line.get("spacing_before", 0))
    return txBox

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # Set corner radius
    shape.adjustments[0] = 0.05
    return shape

def add_slide_number(slide, num, total, color=TEXT_GREY):
    add_text_box(slide, Inches(11.8), Inches(6.9), Inches(1.3), Inches(0.4),
                 f"{num:02d} / {total:02d}", font_size=10, color=color, alignment=PP_ALIGN.RIGHT)

def add_watermark(slide, color=TEXT_GREY):
    add_text_box(slide, Inches(0.5), Inches(6.9), Inches(2), Inches(0.4),
                 "●  Bos Depot Apps", font_size=9, color=color)

def add_image_scaled(slide, img_path, left, top, max_height):
    """Add image maintaining aspect ratio with max height"""
    img = Image.open(img_path)
    w, h = img.size
    ratio = w / h
    height = max_height
    width = int(height * ratio)
    slide.shapes.add_picture(img_path, left, top, width, height)
    return width, height

def add_tag(slide, left, top, text, bg_color, text_color, width=None):
    w = width or Inches(2.2)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, Inches(0.32))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    shape.adjustments[0] = 0.5
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    return shape

def add_bullet_item(slide, left, top, width, title, desc, bullet_color=GOLD, idx=0):
    """Add a bullet point with title and description"""
    h = Inches(0.55)
    y = top + idx * h
    # Bullet dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, y + Inches(0.08), Inches(0.1), Inches(0.1))
    dot.fill.solid()
    dot.fill.fore_color.rgb = bullet_color
    dot.line.fill.background()
    # Text
    txBox = slide.shapes.add_textbox(left + Inches(0.2), y, width - Inches(0.2), h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = title + " — "
    run1.font.size = Pt(12)
    run1.font.color.rgb = TEXT_DARK
    run1.font.bold = True
    run1.font.name = FONT
    run2 = p.add_run()
    run2.text = desc
    run2.font.size = Pt(12)
    run2.font.color.rgb = TEXT_GREY
    run2.font.name = FONT

def add_bullet_item_white(slide, left, top, width, title, desc, idx=0):
    h = Inches(0.55)
    y = top + idx * h
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, y + Inches(0.08), Inches(0.1), Inches(0.1))
    dot.fill.solid()
    dot.fill.fore_color.rgb = GOLD
    dot.line.fill.background()
    txBox = slide.shapes.add_textbox(left + Inches(0.2), y, width - Inches(0.2), h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = title + " — "
    run1.font.size = Pt(12)
    run1.font.color.rgb = WHITE
    run1.font.bold = True
    run1.font.name = FONT
    run2 = p.add_run()
    run2.text = desc
    run2.font.size = Pt(12)
    run2.font.color.rgb = RGBColor(0xBB, 0xCC, 0xDD)
    run2.font.name = FONT

def add_feature_card(slide, left, top, width, height, emoji, title, desc, bg=WHITE, border=RGBColor(0xE2,0xE8,0xF0)):
    card = add_rounded_rect(slide, left, top, width, height, bg, border)
    # Emoji
    add_text_box(slide, left + Inches(0.15), top + Inches(0.12), Inches(0.5), Inches(0.4),
                 emoji, font_size=18, color=TEXT_DARK)
    # Title
    add_text_box(slide, left + Inches(0.65), top + Inches(0.1), width - Inches(0.8), Inches(0.3),
                 title, font_size=11, color=TEXT_DARK, bold=True)
    # Desc
    add_text_box(slide, left + Inches(0.65), top + Inches(0.38), width - Inches(0.8), Inches(0.4),
                 desc, font_size=9, color=TEXT_GREY)

def add_stat_card(slide, left, top, number, label, bg=RGBColor(0x10,0x1E,0x38)):
    w, h = Inches(1.9), Inches(1.2)
    card = add_rounded_rect(slide, left, top, w, h, bg)
    add_text_box(slide, left, top + Inches(0.15), w, Inches(0.5),
                 number, font_size=28, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + Inches(0.7), w, Inches(0.3),
                 label, font_size=10, color=RGBColor(0x88,0x99,0xAA), alignment=PP_ALIGN.CENTER)

def add_impact_card(slide, left, top, emoji, title, desc, bg=WHITE, border=RGBColor(0xE2,0xE8,0xF0)):
    w, h = Inches(2.8), Inches(1.5)
    card = add_rounded_rect(slide, left, top, w, h, bg, border)
    add_text_box(slide, left, top + Inches(0.1), w, Inches(0.4),
                 emoji, font_size=24, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + Inches(0.5), w, Inches(0.35),
                 title, font_size=13, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + Inches(0.85), w, Inches(0.5),
                 desc, font_size=9, color=TEXT_GREY, alignment=PP_ALIGN.CENTER)

def add_tech_pill(slide, left, top, text, dot_color):
    w, h = Inches(1.8), Inches(0.38)
    card = add_rounded_rect(slide, left, top, w, h, WHITE, RGBColor(0xE2,0xE8,0xF0))
    # Dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.12), top + Inches(0.12), Inches(0.12), Inches(0.12))
    dot.fill.solid()
    dot.fill.fore_color.rgb = dot_color
    dot.line.fill.background()
    # Text
    add_text_box(slide, left + Inches(0.3), top + Inches(0.04), w - Inches(0.4), h,
                 text, font_size=10, color=TEXT_DARK, bold=True)


def add_deco_circle(slide, left, top, size, color=WHITE, alpha=0.03):
    """Decorative background circle"""
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    try:
        from pptx.oxml.ns import qn
        srgbClr = circle.element.spPr.find(qn('a:solidFill')).find(qn('a:srgbClr'))
        if srgbClr is not None:
            alpha_el = srgbClr.makeelement(qn('a:alpha'), {})
            alpha_el.set('val', str(int(alpha * 100000)))
            srgbClr.append(alpha_el)
    except Exception:
        pass


# ─── MAIN ───
def build_ppt():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]  # Blank
    total = 12

    img = lambda n: os.path.join(IMG_DIR, f"{n}.png")

    # ══════════════════════════════════════════
    # SLIDE 1: COVER
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_gradient(slide, NAVY, NAVY_LIGHT)
    add_deco_circle(slide, Inches(9), -Inches(2.5), Inches(6))
    add_deco_circle(slide, -Inches(0.5), Inches(5), Inches(4))

    # Left content
    add_tag(slide, Inches(0.8), Inches(1.5), "PORTFOLIO PROJECT", RGBColor(0x1A,0x36,0x57), WHITE, Inches(2))
    add_text_box(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(1.2),
                 "Bos Depot", font_size=48, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(2.75), Inches(5.5), Inches(0.8),
                 "Apps", font_size=48, color=GOLD, bold=True)
    add_text_box(slide, Inches(0.8), Inches(3.6), Inches(5), Inches(0.8),
                 "Sistem Kasir Pintar & Manajemen Gudang untuk\nIndustri Bahan Bangunan — dibangun dengan Flutter & Firebase.",
                 font_size=14, color=RGBColor(0xAA, 0xBB, 0xCC))

    # Developer info
    dev_bg = add_rounded_rect(slide, Inches(0.8), Inches(4.7), Inches(3.3), Inches(0.6),
                               RGBColor(0x10,0x1E,0x38))
    add_text_box(slide, Inches(1.0), Inches(4.75), Inches(3), Inches(0.25),
                 "👨‍💻  Akbar Alfaidah", font_size=12, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.35), Inches(5.0), Inches(3), Inches(0.25),
                 "Mobile Developer — Flutter", font_size=9, color=RGBColor(0x77,0x88,0x99))

    # Right mockups
    add_image_scaled(slide, img(1), Inches(7.2), Inches(0.8), Inches(5.9))
    add_image_scaled(slide, img(2), Inches(9.8), Inches(0.8), Inches(5.9))

    add_slide_number(slide, 1, total, RGBColor(0x55,0x66,0x77))
    add_watermark(slide, RGBColor(0x44,0x55,0x66))

    # ══════════════════════════════════════════
    # SLIDE 2: PROBLEM
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_gradient(slide, WHITE, SURFACE)

    add_tag(slide, Inches(0.8), Inches(1.0), "LATAR BELAKANG", NAVY, GOLD, Inches(2))
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
                 "Masalah yang", font_size=36, color=NAVY, bold=True)
    add_text_box(slide, Inches(0.8), Inches(2.05), Inches(5.5), Inches(0.5),
                 "Ingin Dipecahkan", font_size=36, color=NAVY, bold=True)
    add_text_box(slide, Inches(0.8), Inches(2.7), Inches(5), Inches(0.5),
                 "Usaha panglong kayu & toko bangunan di Indonesia masih dikelola secara tradisional.",
                 font_size=13, color=TEXT_GREY)

    bullets_data = [
        ("Pencatatan Manual", "Buku tulis rawan hilang & salah hitung"),
        ("Tanpa Laporan", "Tidak ada laporan keuangan terstruktur"),
        ("Nota Tulis Tangan", "Tidak profesional & tidak terdokumentasi"),
        ("Tanpa Kontrol Stok", "Sering kehabisan barang tanpa sadar"),
    ]
    for i, (t, d) in enumerate(bullets_data):
        add_bullet_item(slide, Inches(0.8), Inches(3.4), Inches(5.5), t, d, RED, i)

    # Right card
    card_bg = add_rounded_rect(slide, Inches(7.5), Inches(1.2), Inches(4.8), Inches(5),
                                NAVY)
    add_text_box(slide, Inches(7.8), Inches(1.5), Inches(4), Inches(0.3),
                 "DAMPAK MASALAH", font_size=11, color=GOLD, bold=True)

    dampak = [
        ("📉", "Kerugian Tak Terdeteksi", "Tidak tahu mana produk rugi / untung"),
        ("⏰", "Proses Lambat", "Hitung manual per transaksi > 5 menit"),
        ("💸", "Piutang Hilang", "Lupa catat hutang pelanggan"),
        ("📊", "Tanpa Evaluasi", "Tidak bisa analisa tren penjualan"),
    ]
    for i, (emoji, title, desc) in enumerate(dampak):
        y = Inches(2.2) + i * Inches(1.05)
        add_text_box(slide, Inches(7.8), y, Inches(0.5), Inches(0.4), emoji, font_size=22)
        add_text_box(slide, Inches(8.4), y, Inches(3.5), Inches(0.3), title, font_size=13, color=WHITE, bold=True)
        add_text_box(slide, Inches(8.4), y + Inches(0.3), Inches(3.5), Inches(0.3), desc, font_size=10, color=RGBColor(0x88,0x99,0xAA))

    add_slide_number(slide, 2, total)
    add_watermark(slide)

    # ══════════════════════════════════════════
    # SLIDE 3: SOLUTION
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_tag(slide, Inches(0.8), Inches(0.8), "SOLUSI", RGBColor(0xFD,0xF0,0xD5), RGBColor(0xD4,0x85,0x0A), Inches(1.2))
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5),
                 "Satu Aplikasi untuk Semua", font_size=34, color=NAVY, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.9), Inches(5), Inches(0.5),
                 "Bos Depot Apps mengintegrasikan seluruh operasional toko bangunan dalam satu genggaman.",
                 font_size=13, color=TEXT_GREY)

    features = [
        ("🧾", "Smart POS", "Kasir digital dengan negosiasi\nharga & alert keuntungan/kerugian"),
        ("📦", "Gudang Digital", "Stok otomatis terpotong saat\ntransaksi, barcode scanner built-in"),
        ("📊", "Analitik Bisnis", "Dashboard profit real-time, grafik\ntren, top produk terlaris"),
        ("🖨️", "Cetak Nota", "Thermal printer 80mm via\nBluetooth + share ke WhatsApp"),
    ]
    for i, (emoji, title, desc) in enumerate(features):
        col = i % 2
        row = i // 2
        x = Inches(0.8) + col * Inches(2.9)
        y = Inches(2.7) + row * Inches(1.15)
        add_feature_card(slide, x, y, Inches(2.7), Inches(0.95), emoji, title, desc)

    add_image_scaled(slide, img(2), Inches(8.5), Inches(0.6), Inches(6.3))
    add_slide_number(slide, 3, total)
    add_watermark(slide)

    # ══════════════════════════════════════════
    # SLIDE 4: KASIR
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_gradient(slide, WHITE, SURFACE)

    # Left mockup
    add_image_scaled(slide, img(7), Inches(1), Inches(0.7), Inches(6.2))

    # Right content
    add_tag(slide, Inches(6.5), Inches(1.2), "FITUR UTAMA #1", TEAL_BG, TEAL, Inches(1.8))
    add_text_box(slide, Inches(6.5), Inches(1.7), Inches(6), Inches(0.5),
                 "Kasir Pintar dengan Negosiasi", font_size=32, color=NAVY, bold=True)
    add_text_box(slide, Inches(6.5), Inches(2.35), Inches(5.5), Inches(0.6),
                 "Sistem kasir yang memahami dinamika harga di toko bangunan — dimana tawar-menawar adalah hal biasa.",
                 font_size=13, color=TEXT_GREY)

    kasir_bullets = [
        ("Negosiasi Real-time", "Ubah harga per item, sistem tampilkan estimasi untung/rugi"),
        ("Multi-Pembayaran", "Tunai, Transfer, QRIS, dan sistem Hutang"),
        ("Dual Pricing", "Harga eceran (Pcs) dan grosir (Dus, Kodi, Kubik)"),
        ("Auto-Stok", "Stok gudang otomatis berkurang saat checkout"),
    ]
    for i, (t, d) in enumerate(kasir_bullets):
        add_bullet_item(slide, Inches(6.5), Inches(3.2), Inches(6), t, d, GOLD, i)

    add_slide_number(slide, 4, total)
    add_watermark(slide)

    # ══════════════════════════════════════════
    # SLIDE 5: GUDANG
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_tag(slide, Inches(0.8), Inches(0.8), "FITUR UTAMA #2", RGBColor(0xFD,0xF0,0xD5), RGBColor(0xD4,0x85,0x0A), Inches(1.8))
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5),
                 "Manajemen Gudang Terintegrasi", font_size=32, color=NAVY, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.95), Inches(5), Inches(0.6),
                 "Kelola ratusan produk kayu & bangunan dengan sistem yang dirancang khusus untuk industri ini.",
                 font_size=13, color=TEXT_GREY)

    gudang_bullets = [
        ("Kayu & Bangunan", "Form input terpisah: kayu (kelas, dimensi, kubik) & bangunan (Pcs, Kg, Dus)"),
        ("Barcode System", "Generate, cetak label, dan scan barcode untuk percepat input kasir"),
        ("Smart Search", "Pencarian fuzzy yang toleran terhadap typo"),
        ("Drag & Drop", "Atur urutan produk sesuai prioritas toko"),
    ]
    for i, (t, d) in enumerate(gudang_bullets):
        add_bullet_item(slide, Inches(0.8), Inches(2.8), Inches(5.5), t, d, GOLD, i)

    # Right mockups
    add_image_scaled(slide, img(12), Inches(7), Inches(0.6), Inches(6.3))
    add_image_scaled(slide, img(9), Inches(9.8), Inches(0.6), Inches(6.3))

    add_slide_number(slide, 5, total)
    add_watermark(slide)

    # ══════════════════════════════════════════
    # SLIDE 6: LAPORAN & ANALITIK
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_gradient(slide, NAVY, NAVY_LIGHT)
    add_deco_circle(slide, Inches(9), Inches(4), Inches(6))

    # Left mockups
    add_image_scaled(slide, img(6), Inches(0.6), Inches(0.6), Inches(6.3))
    add_image_scaled(slide, img(8), Inches(3.4), Inches(0.6), Inches(6.3))

    # Right content
    add_tag(slide, Inches(7), Inches(1.2), "FITUR UTAMA #3", RGBColor(0x1A,0x36,0x57), WHITE, Inches(1.8))
    add_text_box(slide, Inches(7), Inches(1.7), Inches(5.5), Inches(0.5),
                 "Laporan & Analitik", font_size=32, color=WHITE, bold=True)
    add_text_box(slide, Inches(7), Inches(2.2), Inches(5.5), Inches(0.5),
                 "Bisnis", font_size=32, color=GOLD, bold=True)
    add_text_box(slide, Inches(7), Inches(2.8), Inches(5), Inches(0.5),
                 "Pemilik toko bisa memantau kesehatan bisnis tanpa perlu buka laptop.",
                 font_size=13, color=RGBColor(0xAA,0xBB,0xCC))

    laporan_bullets = [
        ("Top 5 Produk", "Analisa produk terlaris dengan pie chart interaktif"),
        ("Rekap Pendapatan", "Breakdown omset, modal, bensin, dan profit bersih"),
        ("Export CSV", "Download laporan ke Excel dalam satu ketukan"),
        ("Visualisasi Alokasi", "Donut chart pembagian omset secara visual"),
    ]
    for i, (t, d) in enumerate(laporan_bullets):
        add_bullet_item_white(slide, Inches(7), Inches(3.6), Inches(5.5), t, d, i)

    add_slide_number(slide, 6, total, RGBColor(0x55,0x66,0x77))
    add_watermark(slide, RGBColor(0x44,0x55,0x66))

    # ══════════════════════════════════════════
    # SLIDE 7: NOTA & CETAK
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_tag(slide, Inches(0.8), Inches(0.8), "FITUR UTAMA #4", NAVY, GOLD, Inches(1.8))
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5),
                 "Nota Digital & Cetak Thermal", font_size=32, color=NAVY, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.95), Inches(5), Inches(0.6),
                 "Nota profesional yang bisa dicetak langsung ke printer thermal 80mm atau dibagikan via WhatsApp.",
                 font_size=13, color=TEXT_GREY)

    nota_features = [
        ("🖨️", "Cetak Bluetooth", "Support printer Android\n& iOS via Bluetooth"),
        ("📲", "Share WhatsApp", "Kirim nota langsung ke\npelanggan via WA"),
        ("🔄", "Cetak Ulang", "Akses riwayat & cetak\nulang nota kapan saja"),
        ("📸", "Bukti Pembayaran", "Upload foto bukti\ntransfer/pembayaran"),
    ]
    for i, (emoji, title, desc) in enumerate(nota_features):
        col = i % 2
        row = i // 2
        x = Inches(0.8) + col * Inches(2.9)
        y = Inches(2.8) + row * Inches(1.15)
        add_feature_card(slide, x, y, Inches(2.7), Inches(0.95), emoji, title, desc)

    # Right mockups
    add_image_scaled(slide, img(14), Inches(7), Inches(0.6), Inches(6.3))
    add_image_scaled(slide, img(15), Inches(9.8), Inches(0.6), Inches(6.3))

    add_slide_number(slide, 7, total)
    add_watermark(slide)

    # ══════════════════════════════════════════
    # SLIDE 8: ARUS KAS & PIUTANG
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_gradient(slide, WHITE, SURFACE)

    # Left mockups (3 images)
    add_image_scaled(slide, img(3), Inches(0.3), Inches(0.7), Inches(6.2))
    add_image_scaled(slide, img(13), Inches(2.8), Inches(0.7), Inches(6.2))
    add_image_scaled(slide, img(5), Inches(5.3), Inches(0.7), Inches(6.2))

    # Right content
    add_tag(slide, Inches(8.2), Inches(1.2), "FITUR UTAMA #5", TEAL_BG, TEAL, Inches(1.8))
    add_text_box(slide, Inches(8.2), Inches(1.7), Inches(4.5), Inches(0.5),
                 "Arus Kas &", font_size=32, color=NAVY, bold=True)
    add_text_box(slide, Inches(8.2), Inches(2.2), Inches(4.5), Inches(0.5),
                 "Buku Piutang", font_size=32, color=NAVY, bold=True)
    add_text_box(slide, Inches(8.2), Inches(2.8), Inches(4.5), Inches(0.5),
                 "Pantau setiap rupiah yang masuk dan keluar secara real-time.",
                 font_size=13, color=TEXT_GREY)

    arus_bullets = [
        ("Profit Tracker", "Riwayat pemasukan vs pengeluaran + sisa profit bersih"),
        ("Cash Flow", "Timeline setiap transaksi: stok masuk, penjualan, bensin"),
        ("Buku Piutang", "Track hutang per pelanggan + potensi profit piutang"),
        ("Filter Waktu", "Harian, kemarin, 7 hari, bulanan, atau semua waktu"),
    ]
    for i, (t, d) in enumerate(arus_bullets):
        add_bullet_item(slide, Inches(8.2), Inches(3.5), Inches(4.5), t, d, GOLD, i)

    add_slide_number(slide, 8, total)
    add_watermark(slide)

    # ══════════════════════════════════════════
    # SLIDE 9: RIWAYAT TRANSAKSI
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_tag(slide, Inches(0.8), Inches(0.8), "FITUR UTAMA #6", NAVY, GOLD, Inches(1.8))
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5),
                 "Riwayat Transaksi Lengkap", font_size=32, color=NAVY, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.95), Inches(5), Inches(0.6),
                 "Akses seluruh histori penjualan, stok masuk, dan operasional kapan saja dengan filter waktu fleksibel.",
                 font_size=13, color=TEXT_GREY)

    riwayat_bullets = [
        ("Histori Omset", "Semua transaksi tercatat lengkap — status Lunas & Hutang"),
        ("214+ Nota", "Total nota terproses dengan total omset Rp 188 juta+"),
        ("Filter Cerdas", "Hari ini, bulan ini, semua waktu, atau pilih tanggal custom"),
        ("Invoice Detail", "Klik nota untuk lihat detail item, kasir, dan pembayaran"),
    ]
    for i, (t, d) in enumerate(riwayat_bullets):
        add_bullet_item(slide, Inches(0.8), Inches(2.8), Inches(5.5), t, d, GOLD, i)

    add_image_scaled(slide, img(4), Inches(7), Inches(0.6), Inches(6.3))
    add_image_scaled(slide, img(11), Inches(9.8), Inches(0.6), Inches(6.3))

    add_slide_number(slide, 9, total)
    add_watermark(slide)

    # ══════════════════════════════════════════
    # SLIDE 10: KEAMANAN
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_gradient(slide, WHITE, SURFACE)

    add_tag(slide, Inches(0.8), Inches(0.8), "KEAMANAN", NAVY, GOLD, Inches(1.5))
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5),
                 "Multi-Layer Security System", font_size=32, color=NAVY, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.95), Inches(5), Inches(0.5),
                 "6 lapisan keamanan melindungi data bisnis dari akses tidak sah.",
                 font_size=13, color=TEXT_GREY)

    security = [
        ("🔐", "Firebase Auth", "Login Email & Google OAuth 2.0"),
        ("🔑", "PIN Keamanan", "PIN 6 digit khusus pemilik toko"),
        ("👆", "Biometrik", "Sidik jari & Face ID"),
        ("⏱️", "Auto-Logout", "Sesi pemilik expired 10 menit"),
        ("🛡️", "Enkripsi Sesi", "Secure storage terenkripsi"),
        ("👥", "Role-Based", "Akses berbeda: Owner vs Karyawan"),
    ]
    for i, (emoji, title, desc) in enumerate(security):
        col = i % 2
        row = i // 2
        x = Inches(0.8) + col * Inches(2.9)
        y = Inches(2.7) + row * Inches(1.15)
        add_feature_card(slide, x, y, Inches(2.7), Inches(0.95), emoji, title, desc)

    add_image_scaled(slide, img(1), Inches(8.5), Inches(0.6), Inches(6.3))

    add_slide_number(slide, 10, total)
    add_watermark(slide)

    # ══════════════════════════════════════════
    # SLIDE 11: TECH STACK
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_gradient(slide, NAVY, NAVY_LIGHT)
    add_deco_circle(slide, -Inches(1), -Inches(2), Inches(5))

    add_tag(slide, Inches(4.6), Inches(0.6), "TEKNOLOGI", RGBColor(0x1A,0x36,0x57), WHITE, Inches(1.6))
    add_text_box(slide, Inches(1), Inches(1.1), Inches(11), Inches(0.5),
                 "Tech Stack yang Digunakan", font_size=34, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2), Inches(1.7), Inches(9), Inches(0.5),
                 "Dibangun dengan teknologi modern untuk performa, skalabilitas, dan pengalaman pengguna terbaik.",
                 font_size=13, color=RGBColor(0xAA,0xBB,0xCC), alignment=PP_ALIGN.CENTER)

    tech_items = [
        ("Flutter & Dart", RGBColor(0x00,0x78,0xD4)),
        ("Firebase Firestore", RGBColor(0xFF,0x6F,0x00)),
        ("Firebase Auth", RGBColor(0xEA,0x43,0x35)),
        ("SQLite (Offline)", GREEN),
        ("GetX", RGBColor(0x7C,0x3A,0xED)),
        ("ESC/POS Printing", TEAL),
        ("Google Sign-In", RGBColor(0x00,0x78,0xD4)),
        ("FL Chart", RGBColor(0xFF,0x6F,0x00)),
        ("Bluetooth Thermal", GREEN),
        ("Biometric Auth", RGBColor(0xEA,0x43,0x35)),
        ("CSV Export", RGBColor(0x7C,0x3A,0xED)),
        ("Barcode Scanner", TEAL),
    ]
    for i, (name, dot_c) in enumerate(tech_items):
        col = i % 6
        row = i // 6
        x = Inches(0.7) + col * Inches(2.05)
        y = Inches(2.5) + row * Inches(0.55)
        add_tech_pill(slide, x, y, name, dot_c)

    # Stats
    stats = [("33", "Halaman UI"), ("22", "Controller"), ("24", "Data Source"), ("30+", "Dependencies")]
    for i, (num, label) in enumerate(stats):
        add_stat_card(slide, Inches(1.4) + i * Inches(2.6), Inches(4.2), num, label)

    add_slide_number(slide, 11, total, RGBColor(0x55,0x66,0x77))
    add_watermark(slide, RGBColor(0x44,0x55,0x66))

    # ══════════════════════════════════════════
    # SLIDE 12: DAMPAK & CLOSING
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_gradient(slide, WHITE, SURFACE)

    add_tag(slide, Inches(4.5), Inches(0.4), "DAMPAK & HASIL", NAVY, GOLD, Inches(1.8))
    add_text_box(slide, Inches(1), Inches(0.85), Inches(11), Inches(0.5),
                 "Dampak Nyata untuk Bisnis Klien", font_size=34, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2), Inches(1.45), Inches(9), Inches(0.5),
                 "Aplikasi ini sudah berjalan dan digunakan oleh klien toko bangunan secara aktif setiap hari.",
                 font_size=13, color=TEXT_GREY, alignment=PP_ALIGN.CENTER)

    impacts = [
        ("💰", "Rp 188 Juta+", "Total omset yang berhasil dikelola\nmelalui sistem aplikasi"),
        ("🧾", "214+ Nota Terproses", "Transaksi tercatat rapi dengan\ninvoice otomatis"),
        ("📋", "Rp 170 Juta Piutang", "Tercatat & terlacak — tidak ada\nlagi hutang yang terlupakan"),
        ("⚡", "10x Lebih Cepat", "Proses transaksi dari 5 menit\nmenjadi 30 detik per nota"),
        ("📱", "Cross-Platform", "Berjalan di Android & iOS\ndengan satu codebase Flutter"),
        ("☁️", "Cloud-Synced", "Data aman di cloud Firebase\n+ bisa diakses offline"),
    ]
    for i, (emoji, title, desc) in enumerate(impacts):
        col = i % 3
        row = i // 3
        x = Inches(1.3) + col * Inches(3.7)
        y = Inches(2.2) + row * Inches(1.7)
        add_impact_card(slide, x, y, emoji, title, desc)

    # CTA
    cta = add_rounded_rect(slide, Inches(4.2), Inches(6.0), Inches(4.8), Inches(0.55), GOLD)
    add_text_box(slide, Inches(4.2), Inches(6.05), Inches(4.8), Inches(0.45),
                 "💼  Hubungi Saya di LinkedIn", font_size=13, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2), Inches(6.65), Inches(9), Inches(0.3),
                 "Akbar Alfaidah — Mobile Developer  •  Flutter & Firebase Specialist",
                 font_size=10, color=TEXT_GREY, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 12, total)
    add_watermark(slide)

    # ─── SAVE ───
    prs.save(OUT_FILE)
    print(f"\n✅ PPT berhasil di-generate!")
    print(f"📂 File: {OUT_FILE}")
    print(f"📊 Total: {total} slide")


if __name__ == "__main__":
    build_ppt()
